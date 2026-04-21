"""
Presentation Export Utilities - Create PPTX from screenshots.

Educational Note: This module uses python-pptx to create PowerPoint presentations
from PNG screenshots captured from HTML slides.

Why python-pptx?
- Pure Python library (no external dependencies)
- Well-documented and stable
- Creates .pptx files compatible with PowerPoint, Google Slides, etc.

Workflow:
1. Screenshots are captured at 1920x1080 (16:9 aspect ratio)
2. Each screenshot becomes a full-slide image in the PPTX
3. Slide dimensions are set to standard widescreen (13.333" x 7.5")

Usage:
    from app.utils.presentation_export_utils import create_pptx_from_screenshots

    pptx_path = create_pptx_from_screenshots(
        screenshots=[{"screenshot_path": "/path/to/slide_01.png", ...}],
        output_path="/path/to/presentation.pptx",
        title="My Presentation"
    )
"""

import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)


# Standard 16:9 widescreen dimensions (in inches)
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5


def create_pptx_from_screenshots(
    screenshots: List[Dict[str, Any]],
    output_path: str,
    title: str = "Presentation",
    progress_callback: Optional[callable] = None
) -> Optional[str]:
    """
    Create a PPTX file from screenshot images.

    Educational Note: Each screenshot becomes a full-slide background image.
    The PPTX maintains the exact visual design from the HTML slides.

    Args:
        screenshots: List of screenshot dicts from capture_slides_as_screenshots
            Each dict should have:
            - screenshot_path: Path to the PNG file
            - success: Boolean indicating if screenshot was captured
        output_path: Path to save the PPTX file
        title: Presentation title (for metadata)
        progress_callback: Optional callback for progress updates

    Returns:
        Path to the created PPTX file, or None if failed

    Example:
        pptx_path = create_pptx_from_screenshots(
            screenshots=[
                {"screenshot_path": "/path/slide_01.png", "success": True},
                {"screenshot_path": "/path/slide_02.png", "success": True},
            ],
            output_path="/path/to/output.pptx",
            title="Q4 Report"
        )
    """
    try:
        # Filter to successful screenshots
        valid_screenshots = [s for s in screenshots if s.get("success") and s.get("screenshot_path")]

        if not valid_screenshots:
            logger.warning("No valid screenshots to export")
            return None

        # Create presentation with widescreen dimensions
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
        prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

        # Set presentation properties
        prs.core_properties.title = title

        # Get blank layout (no placeholders)
        blank_layout = prs.slide_layouts[6]  # Blank layout

        for i, screenshot in enumerate(valid_screenshots):
            screenshot_path = screenshot["screenshot_path"]

            if progress_callback:
                progress_callback("adding_slide", {
                    "slide": screenshot.get("slide_file", f"Slide {i + 1}"),
                    "index": i + 1,
                    "total": len(valid_screenshots)
                })

            # Verify screenshot file exists
            if not Path(screenshot_path).exists():
                logger.warning("Screenshot not found, skipping: %s", screenshot_path)
                continue

            # Add blank slide
            slide = prs.slides.add_slide(blank_layout)

            # Add screenshot as full-slide background image
            # Position at (0, 0) with full slide dimensions
            slide.shapes.add_picture(
                screenshot_path,
                Inches(0),
                Inches(0),
                width=Inches(SLIDE_WIDTH_INCHES),
                height=Inches(SLIDE_HEIGHT_INCHES)
            )

        # Ensure output directory exists
        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)

        # Save the presentation
        prs.save(str(output_file))

        logger.info("Created presentation with %s slides: %s", len(valid_screenshots), output_path)

        return str(output_file)

    except Exception as e:
        logger.exception("Error creating presentation")
        return None


def get_slide_count(pptx_path: str) -> int:
    """
    Get the number of slides in a PPTX file.

    Args:
        pptx_path: Path to the PPTX file

    Returns:
        Number of slides, or 0 if file doesn't exist/is invalid
    """
    try:
        prs = Presentation(pptx_path)
        return len(prs.slides)
    except Exception:
        return 0


def is_pptx_available() -> bool:
    """
    Check if python-pptx is installed.

    Returns:
        True if python-pptx can be used, False otherwise
    """
    try:
        from pptx import Presentation
        return True
    except ImportError:
        return False
